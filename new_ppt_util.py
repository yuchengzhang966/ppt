#!/usr/bin/env python
# coding: utf-8

from pptx import Presentation
from pptx.util import Inches, Pt
import os
import matplotlib.pyplot as plt
import matplotlib.image as mpimg
from PIL import Image
import csv
import fitz  # PyMuPDF
import PyPDF2

def add_title_slide(prs, title_text, subtitle_text):
    """
    Adds a title slide to the presentation with the specified title and subtitle.

    Parameters:
    prs (Presentation): The PowerPoint presentation object.
    title_text (str): The text for the slide's title.
    subtitle_text (str): The text for the slide's subtitle.
    """
    # Select the title slide layout
    fengmian = prs.slide_layouts[0]
    
    # Add a new slide with the selected layout
    slide = prs.slides.add_slide(fengmian)
    
    # Set the title and subtitle
    title = slide.shapes.title
    subtitle = slide.placeholders[1]

    title.text = title_text
    subtitle.text = subtitle_text

def add_title_only_slide(prs, title_text):
    """
    Adds a title-only slide to the presentation with the specified title.

    Parameters:
    prs (Presentation): The PowerPoint presentation object.
    title_text (str): The text for the slide's title.
    """
    # Select the title-only slide layout
    biaoti = prs.slide_layouts[2]
    
    # Add a new slide with the selected layout
    slide = prs.slides.add_slide(biaoti)
    
    # Set the title
    title = slide.shapes.title
    title.text = title_text

def add_image_slide(prs, img_path, title_text=""):
    """
    Adds a slide with an image that is resized to fit within the slide's dimensions,
    centered horizontally and vertically.

    Parameters:
    prs (Presentation): The PowerPoint presentation object.
    img_path (str): The path to the image file.
    title_text (str): The text for the slide's title. Default is an empty string.
    """
    # Load the image using Pillow to get its dimensions
    image = Image.open(img_path)
    img_width, img_height = image.size
    
    # Get the slide dimensions
    slide_width = prs.slide_width - Pt(180)
    slide_height = prs.slide_height - Pt(150)
    
    # Determine the aspect ratios
    image_aspect_ratio = img_width / img_height
    slide_aspect_ratio = slide_width / slide_height
    
    # Decide whether to fit the image to the slide's width or height
    if image_aspect_ratio > slide_aspect_ratio:
        # Image is wider than the slide; fit to width
        new_width = slide_width
        new_height = new_width / image_aspect_ratio
    else:
        # Image is taller than the slide; fit to height
        new_height = slide_height 
        new_width = new_height * image_aspect_ratio  
    
    # Center the image horizontally and move it down by 100px
    left = (slide_width - new_width) / 2 + Pt(100)
    top = (slide_height - new_height) / 2 + Pt(90)  # Move down by 100px
    
    # Add a slide layout and insert the image with the calculated dimensions
    tu = prs.slide_layouts[3]
    slide = prs.slides.add_slide(tu)
    
    # Set the title for the slide
    title = slide.shapes.title
    title.text = title_text
    
    # Insert the image with the calculated size and position
    pic = slide.shapes.add_picture(img_path, left, top, width=new_width, height=new_height)

def generate_presentation(csv_file_path, pr_path, output_path, folder_path):
    """
    Generates a PowerPoint presentation based on data from a CSV file.

    Parameters:
    csv_file_path (str): The path to the CSV file containing presentation data.
    pr_path (str): The path to the PowerPoint template.
    output_path (str): The path to save the generated presentation.
    folder_path (str): The path to the folder containing image files.
    """
    # Load the PowerPoint template
    prs = Presentation(pr_path)

    with open(csv_file_path, mode='r') as file:
        # Create a CSV reader object
        csv_reader = csv.reader(file)
        
        # Skip the header row if there is one
        next(csv_reader, None)
        
        # Iterate over each row in the CSV file
        for row in csv_reader:
            number, title, stype, img, subtitle = row
            
            # Call the appropriate function based on the Type column
            if stype == "title page":
                add_title_slide(prs, title, subtitle_text=subtitle)
            elif stype == "table of content":
                add_content_slide(prs, title, subtitle_text=subtitle)
            elif stype == "subtitle":
                add_title_only_slide(prs, title)
            elif stype == "main":
                add_image_slide(prs, img_path=os.path.join(folder_path, 'output_images', img), title_text=title)

    # Save the presentation to the specified output path
    prs.save(output_path)
    print("Successfully generated presentation")

def trim_space(image_path, output_path, white_threshold=100):
    """
    Trims the image to a specified box, removes any remaining white space, and then adds a 10-pixel margin.

    Parameters:
    image_path (str): Path to the input image.
    output_path (str): Path to save the trimmed image.
    white_threshold (int): RGB threshold to consider a pixel as white.
    """
    # Open the image
    image = Image.open(image_path).convert("RGB")

    # Get the dimensions of the image
    width, height = image.size

    # First trim the image with the specified box
    box = (5, 100, width - 5, height - 66)
    image = image.crop(box)

    # Convert to a binary image where white pixels are 0 and non-white pixels are 255
    binary_image = image.point(lambda x: 0 if x >= white_threshold else 255)

    # Convert to grayscale to work with getbbox
    binary_image = binary_image.convert("L")

    # Get the bounding box of the non-white area
    bbox = binary_image.getbbox()
    bbox_1 = (0, bbox[1], width - 10, bbox[3])
    bbox = bbox_1

    # Crop the image to the bounding box
    if bbox:
        final_trimmed_image = image.crop(bbox)
        final_trimmed_image.save(output_path)
    else:
        print("No non-white area found. The image was not cropped.")

def add_content_slide(prs, title_text, subtitle_text):
    """
    Adds a content slide to the presentation with the specified title and subtitle.

    Parameters:
    prs (Presentation): The PowerPoint presentation object.
    title_text (str): The text for the slide's title.
    subtitle_text (str): The text for the slide's subtitle (can be multiline).
    """
    # Select the content slide layout
    mulu = prs.slide_layouts[1]
    
    # Add a new slide with the selected layout
    slide = prs.slides.add_slide(mulu)
    
    # Set the title and subtitle
    title = slide.shapes.title
    subtitle = slide.placeholders[1]

    title.text = title_text
    subtitle.text = subtitle_text

def extract_text_from_page(pdf_path, page_number):
    """
    Extracts text from a specific page in a PDF.

    Parameters:
    pdf_path (str): Path to the PDF file.
    page_number (int): The page number to extract text from.

    Returns:
    str: Extracted text or None if the page number is out of range.
    """
    # Open the PDF file
    with open(pdf_path, 'rb') as pdf_file:
        # Create a PDF reader object
        pdf_reader = PyPDF2.PdfReader(pdf_file)
        
        # Check if the page number is within the range
        if page_number < 1 or page_number > len(pdf_reader.pages):
            print(f"Error: Page number {page_number} is out of range.")
            return None
        
        # Get the specific page
        page = pdf_reader.pages[page_number - 1]  # Pages are 0-indexed
        
        # Extract the text from the page
        text = page.extract_text()
        
        return text

def pdf_to_png_and_trim(pdf_path, output_folder, dpi=300):
    """
    Converts each page of a PDF to PNG.

    :param pdf_path: Path to the input PDF file.
    :param output_folder: Folder to save the output PNG files.
    :param dpi: Dots per inch for the output PNG.
    """
    # Create the output directory if it doesn't exist
    if not os.path.exists(output_folder):
        os.makedirs(output_folder)

    # Open the PDF file
    pdf_document = fitz.open(pdf_path)

    # Iterate through each page
    for page_number in range(len(pdf_document)):
        page = pdf_document.load_page(page_number)  # Load the page
        zoom = dpi / 72  # Calculate the zoom factor
        mat = fitz.Matrix(zoom, zoom)  # Create the transformation matrix
        pix = page.get_pixmap(matrix=mat)  # Render page to an image
        output_path = os.path.join(output_folder, f"page_{page_number + 1}.png")
        pix.save(output_path)  # Save the image as a PNG file
        trim_space(output_path,output_path)
        print("finished trimming space")



def extract_and_trim_pdf_pages(pdf_path, output_folder, output_images_folder, page_numbers):
    """
    Extracts specified pages from a PDF and trims white space.

    Parameters:
    pdf_path (str): Path to the input PDF file.
    output_folder (str): Path to the folder for saving the output PNG files.
    output_images_folder (str): Path to the folder for saving trimmed images.
    page_numbers (list): List of page numbers to extract and trim.
    """
    # Create output folders if they don't exist
    if not os.path.exists(output_folder):
        os.makedirs(output_folder)
    if not os.path.exists(output_images_folder):
        os.makedirs(output_images_folder)

    # Open the PDF
    pdf_document = fitz.open(pdf_path)

    for page_number in page_numbers:
        page = pdf_document.load_page(page_number - 1)  # Load page (0-indexed)
        zoom = 3  # Set the zoom level
        mat = fitz.Matrix(zoom, zoom)  # Transformation matrix
        pix = page.get_pixmap(matrix=mat)  # Get the pixmap for the page
        image_path = os.path.join(output_folder, f"page_{page_number}.png")
        pix.save(image_path)  # Save the page as a PNG
        print(f"Extracted page {page_number} to {image_path}")
        
        # Trim the image
        trimmed_image_path = os.path.join(output_images_folder, f"trimmed_page_{page_number}.png")
        trim_space(image_path, trimmed_image_path)
        print(f"Trimmed page {page_number} saved to {trimmed_image_path}")

    # Close the PDF document
    pdf_document.close()

# Function to extract slide information from a PDF
def extract_slide_info(pdf_path, output_csv_path,slide_title):
    # Open the PDF file
    with open(pdf_path, 'rb') as pdf_file:
        pdf_reader = PyPDF2.PdfReader(pdf_file)

        # Initialize a list to hold titles and their corresponding type
        slides_info = []

        # Iterate through the pages and extract titles and categorize
        for i, page in enumerate(pdf_reader.pages):
            # Extract text from the page
            text = page.extract_text().strip()
            text = text.encode('utf-8').decode('utf-8')

            # Determine the type of the slide
            if i == 0:
                slide_type = "title page"
                
            elif "Table of Contents" in text:
                slide_type = "table of content"
                subtitle = extract_middle_lines(text)  # Extract the ToC as subtitle
                print(subtitle)
                print("table of contetns")
                print(text)
            else:
                slide_type = "main"
                subtitle = ""

            # Find the title for the slide
            title_lines = text.split('\n')
            if len(title_lines) > 0:
                title = title_lines[0]  # Assume the first line is the title
                if i == 0 :
                    title = slide_title
                    subtitle = title_lines[0]
            else:
                title = "No Title"

            # Generate the image file name based on the page number
            image_name = f"page_{i + 1}.png"

            # Append the slide info to the list, including the subtitle
            slides_info.append([i + 1, title, slide_type, image_name, subtitle])

    # Write the extracted information to a CSV file
    with open(output_csv_path, 'w', newline='', encoding='utf-8') as csv_file:
        writer = csv.writer(csv_file)
        # Write the header
        writer.writerow(["Page Number", "Title", "Slide Type", "Image Name", "Subtitle"])
        # Write the slide information
        writer.writerows(slides_info)
    print("extracted slide info to"+output_csv_path)

# Function to extract the middle lines of text, excluding the first and last lines
def extract_middle_lines(text):
    # Split the text into individual lines
    lines = text.strip().split('\n')
    
    # Check if there are at least three lines (to exclude first and last)
    if len(lines) > 2:
        # Exclude the first and last lines
        middle_lines = lines[1:-1]
    else:
        # If there are fewer than three lines, return an empty string or handle as needed
        return ""

    # Join the remaining lines back into a single string
    result_text = "\n".join(middle_lines)
    return result_text



def insert_row_after_page_number(csv_file_path, page_number, title):
    # Read the CSV file and store the rows in a list
    with open(csv_file_path, mode='r') as file:
        reader = csv.reader(file)
        headers = next(reader)  # Read the headers
        rows = list(reader)

    # Find the index of the row with the specified page number
    page_number_index = headers.index("Page Number")
    title_index = headers.index("Title")
    slide_type_index = headers.index("Slide Type")
    
    for i, row in enumerate(rows):
        if int(row[page_number_index]) == page_number:  # Assuming "Page Number" column exists
            # Insert the new row after the found row
            new_row = [""] * len(headers)  # Create an empty row with the same number of columns
            new_row[page_number_index] = page_number
            new_row[title_index] = title
            new_row[slide_type_index] = "subtitle"
            rows.insert(i + 1, new_row)
            break

    # Write the updated rows back to the CSV file
    with open(csv_file_path, mode='w', newline='') as file:
        writer = csv.writer(file)
        writer.writerow(headers)  # Write the headers back
        writer.writerows(rows)

